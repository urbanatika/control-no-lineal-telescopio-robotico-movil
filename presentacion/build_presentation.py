"""
Genera presentacion.pptx a partir de los resultados de la simulación
(sim_results.pkl) y las figuras en ../figures/. Estilo limpio, consistente,
apto para exponer.
"""
import pickle
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIG = "../figures"

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0xD6, 0x27, 0x28)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF3, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_title(slide, text, size=32, color=NAVY, top=0.4, left=0.6, width=12.1, bold=True):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def add_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_bullets(slide, items, left=0.7, top=1.5, width=11.9, height=5.3, size=18, color=RGBColor(0x22, 0x22, 0x22)):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (level, text) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size - level * 2)
        r.font.color.rgb = color
        r.font.name = "Calibri"
        p.space_after = Pt(8)
    return box


def add_image(slide, path, left, top, width=None, height=None):
    kwargs = {}
    if width:
        kwargs["width"] = Inches(width)
    if height:
        kwargs["height"] = Inches(height)
    return slide.shapes.add_picture(path, Inches(left), Inches(top), **kwargs)


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(12), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.color.rgb = GRAY


# ---------------------------------------------------------------------
# 1. Portada
# ---------------------------------------------------------------------
s = add_slide()
bg(s, NAVY)
box = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.2))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Control No Lineal Multi-Estrategia para un\nTelescopio Robótico Móvil"
r.font.size = Pt(38)
r.font.bold = True
r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "Adaptativo-Lyapunov + EKF vs. SMC vs. MPC vs. Neuro-Difuso ante perturbación sísmica, viento y vibración de base móvil"
r2.font.size = Pt(18)
r2.font.color.rgb = RGBColor(0xC8, 0xD0, 0xE0)
p2.space_before = Pt(16)

box2 = s.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(1.0))
tf2 = box2.text_frame
p = tf2.paragraphs[0]
r = p.add_run()
r.text = "Mag. Diego N. Elorza-Rivera  ·  Jorge Esper\nUniversidad Andrés Bello — Facultad de Ingeniería"
r.font.size = Pt(14)
r.font.color.rgb = RGBColor(0xAA, 0xB4, 0xC8)

# ---------------------------------------------------------------------
# 2. Motivación
# ---------------------------------------------------------------------
s = add_slide(); bg(s); add_bar(s)
add_title(s, "De la antena fija de ALMA al telescopio robótico móvil")
add_bullets(s, [
    (0, "Trabajo base: antena ALMA de 12 m (~100 t) — control Adaptativo-Lyapunov + feedforward sísmico + EKF"),
    (0, "Nuevo desafío: robot móvil + brazo alt-azimutal de 2 GDL sosteniendo un instrumento óptico"),
    (1, "Inercia mucho menor (J ~ 1–3 kg·m²) → misma perturbación produce error angular relativo mucho mayor"),
    (1, "Nuevo canal de perturbación: microvibración de terreno/ruedas (5–50 Hz), sin análogo en el trabajo original"),
    (0, "Pregunta de investigación: ¿qué ley de control resiste mejor esta combinación de perturbaciones?"),
], top=1.7)
add_footer(s, "2")

# ---------------------------------------------------------------------
# 3. Modelo dinámico
# ---------------------------------------------------------------------
s = add_slide(); bg(s); add_bar(s)
add_title(s, "Modelo dinámico: manipulador alt-azimutal de 2 GDL")
add_bullets(s, [
    (0, "Extensión vectorial de Euler-Lagrange (ec. 1 del trabajo base) a acimut (q1) + elevación (q2)"),
    (0, "Acoplamiento inercial J1(q2) = J1,0 + J1,c·cos²(q2)"),
    (0, "Perturbación total por eje:"),
    (1, "τ_s = Fs·ẍg(t)  — sísmico (feedforward vía acelerómetro de fundación)"),
    (1, "τ_v — viento, proceso de Gauss-Markov de 1er orden (idéntico al trabajo base)"),
    (1, "NUEVO: microvibración de terreno — ruido banda 5–50 Hz, NO compensado por feedforward"),
    (0, "Escenario: 60 s, perfil multi-evento — Basal / Coquimbo 2015 / Chiloé 2016 / Melipilla 2017"),
], top=1.6)
add_footer(s, "3")

# ---------------------------------------------------------------------
# 4. Las 4 leyes de control
# ---------------------------------------------------------------------
s = add_slide(); bg(s); add_bar(s)
add_title(s, "Cuatro leyes de control, un EKF de fusión sensorial común")
cols = [
    ("A. Adaptativo-Lyapunov", "Extensión directa del trabajo base.\nLey de adaptación paramétrica\nθ̇ = -Γ Yᵀs. Compensa error\nde modelo por convergencia."),
    ("B. SMC Super-Twisting", "Modos deslizantes de 2° orden\n(Moreno & Osorio 2012). Robusto\nsin adaptación; reduce chattering\nrespecto al SMC clásico."),
    ("C. MPC lineal", "Horizonte deslizante, QP en forma\ncerrada (batch). Restricciones de\npar. Sin corrección de sesgo\nparamétrico (no offset-free)."),
    ("D. Neuro-Difuso", "Base gaussiana tipo Wang (1993).\nAproxima dinámica residual no\nmodelada. Ley de adaptación\nderivada de Lyapunov."),
]
left = 0.6
w = 2.95
for i, (title, body) in enumerate(cols):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + i * (w + 0.15)), Inches(1.7), Inches(w), Inches(4.6))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = NAVY; box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = body
    r2.font.size = Pt(12.5); r2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p2.space_before = Pt(10)
add_footer(s, "4")

# ---------------------------------------------------------------------
# 5. EKF - estimación de viento
# ---------------------------------------------------------------------
s = add_slide(); bg(s); add_bar(s)
add_title(s, "EKF: estimación robusta del torque de viento no medible")
add_image(s, f"{FIG}/fig3_ekf_estimacion_viento.png", left=1.3, top=1.5, width=10.7)
add_footer(s, "5  ·  RMSE estimación: 2.14 N·m (acimut), 2.03 N·m (elevación)")

# ---------------------------------------------------------------------
# 6. Resultados - error comparativo
# ---------------------------------------------------------------------
s = add_slide(); bg(s); add_bar(s)
add_title(s, "Resultado principal: error de apuntamiento comparado")
add_image(s, f"{FIG}/fig1_error_comparativo.png", left=0.9, top=1.4, width=11.5)
add_footer(s, "6")

# ---------------------------------------------------------------------
# 7. RMSE por fase
# ---------------------------------------------------------------------
s = add_slide(); bg(s); add_bar(s)
add_title(s, "RMSE por fase sísmica")
add_image(s, f"{FIG}/fig7_rmse_por_fase.png", left=1.5, top=1.5, width=10.3)
add_footer(s, "7")

# ---------------------------------------------------------------------
# 8. Tabla de métricas
# ---------------------------------------------------------------------
s = add_slide(); bg(s); add_bar(s)
add_title(s, "Métricas comparativas (60 s de simulación)")

with open("../python_sim/sim_results.pkl", "rb") as f:
    data = pickle.load(f)
metrics = data["metrics"]

rows = len(metrics) + 1
cols = 6
tbl_shape = s.shapes.add_table(rows, cols, Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.6))
table = tbl_shape.table
headers = ["Controlador", "RMSE total\n[arcsec]", "Máx. total\n[arcsec]", "Esfuerzo\n∫τ²dt", "Chattering\n[N·m]", "Cómputo/paso\n[ms]"]
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(12)

best_rmse = min(m["rmse_total"] for m in metrics)
for i, m in enumerate(metrics, start=1):
    vals = [m["label"], f"{m['rmse_total']:.1f}", f"{m['max_total']:.1f}",
            f"{m['effort']:.0f}", f"{m['chattering']:.0f}", f"{m['compute_ms']:.4f}"]
    for j, v in enumerate(vals):
        cell = table.cell(i, j)
        cell.text = v
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xE8) if (j == 1 and m["rmse_total"] == best_rmse) else WHITE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(11.5)
                r.font.bold = (j == 1 and m["rmse_total"] == best_rmse)

add_bullets(s, [
    (0, "SMC (Super-Twisting): mejor RMSE total — robustez sin necesidad de convergencia paramétrica"),
    (0, "MPC lineal: mayor error sostenido — no corrige el sesgo de parámetros nominales (no offset-free)"),
    (0, "Las 4 leyes cumplen holgadamente el presupuesto de cómputo de un lazo a 1 kHz"),
], top=4.7, size=16)
add_footer(s, "8")

# ---------------------------------------------------------------------
# 9. Convergencia paramétrica
# ---------------------------------------------------------------------
s = add_slide(); bg(s); add_bar(s)
add_title(s, "Convergencia paramétrica — ley Adaptativa-Lyapunov")
add_image(s, f"{FIG}/fig4_convergencia_parametrica.png", left=0.5, top=2.0, width=12.3)
add_bullets(s, [
    (0, "mgl converge a su valor real en < 20 s (regresor con excitación persistente fuerte, cos(q2))"),
    (0, "J, b convergen mucho más lento — excitación persistente débil bajo referencia casi constante"),
], top=1.15, size=15)
add_footer(s, "9")

# ---------------------------------------------------------------------
# 10. Materiales / prototipo físico
# ---------------------------------------------------------------------
s = add_slide(); bg(s); add_bar(s)
add_title(s, "Hoja de ruta física: prototipo de banco de laboratorio")
add_bullets(s, [
    (0, "Base móvil (motores DC+encoder, ruedas, batería/BMS): ~US$550"),
    (0, "Brazo robótico 2–3 GDL (actuadores, estructura, rodamientos): ~US$558"),
    (0, "Carga útil óptica (láser + PSD/cámara de guiado): ~US$130"),
    (0, "Sensórica (IMU, acelerómetro de base, encoders): ~US$75 (+US$300 opcional FOG)"),
    (0, "Cómputo (MCU tiempo real 1 kHz + SBC): ~US$170"),
    (0, "Banco de pruebas (shaker + ventiladores para sismo/viento sintéticos): ~US$260"),
    (0, "COSTO TOTAL ESTIMADO: ~US$1 743 (base) / ~US$2 043 (con FOG)"),
], top=1.6, size=18)
add_footer(s, "10")

# ---------------------------------------------------------------------
# 11. Arquitectura del prototipo (diagrama de bloques)
# ---------------------------------------------------------------------
s = add_slide(); bg(s); add_bar(s)
add_title(s, "Arquitectura del prototipo: diagrama de bloques")
add_image(s, f"{FIG}/fig9_diagrama_bloques.png", left=3.57, top=1.5, height=5.2)
add_footer(s, "11")

# ---------------------------------------------------------------------
# 12. Layout mecánico + esquema eléctrico
# ---------------------------------------------------------------------
s = add_slide(); bg(s); add_bar(s)
add_title(s, "Diseño mecánico y esquema eléctrico")
add_image(s, f"{FIG}/fig11_layout_mecanico.png", left=0.6, top=1.7, height=4.7)
add_image(s, f"{FIG}/fig10_esquema_electrico.png", left=6.7, top=1.7, height=4.7)
add_footer(s, "12")

# ---------------------------------------------------------------------
# 13. Banco de pruebas
# ---------------------------------------------------------------------
s = add_slide(); bg(s); add_bar(s)
add_title(s, "Banco de pruebas: perturbaciones controladas")
add_image(s, f"{FIG}/fig12_banco_pruebas.png", left=2.59, top=1.7, height=5.0)
add_footer(s, "13")

# ---------------------------------------------------------------------
# 14. Conclusiones
# ---------------------------------------------------------------------
s = add_slide(); bg(s); add_bar(s)
add_title(s, "Conclusiones")
add_bullets(s, [
    (0, "SMC Super-Twisting: mayor robustez global frente a incertidumbre paramétrica sostenida y perturbación no modelada"),
    (0, "Adaptativo-Lyapunov y Neuro-Difuso: requieren convergencia inicial, luego igualan aprox. al SMC"),
    (0, "MPC lineal: limitado por ausencia de corrección de sesgo — candidato a mejora vía esquema offset-free"),
    (0, "EKF de fusión sensorial: estimación robusta del viento, válida como capa común independiente de la ley de control"),
    (0, "Próximo paso: validación física en el prototipo de banco de laboratorio propuesto"),
], top=1.7, size=19)
add_footer(s, "14")

prs.save("presentacion.pptx")
print("Guardado presentacion.pptx")
